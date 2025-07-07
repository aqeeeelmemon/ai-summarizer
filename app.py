import streamlit as st
import openai
import PyPDF2
import docx
import io
from pptx import Presentation
from fpdf import FPDF

# --- Streamlit Config ---
st.set_page_config(page_title="AI File Summarizer", layout="centered")
st.title("📁 AI File Summarizer & Flowchart Generator")
st.markdown("<h4 style='text-align: center; color: gray;'>Made by Aqeel Memon</h4>", unsafe_allow_html=True)

# --- OpenRouter Setup ---
api_key = st.sidebar.text_input("🔑 Enter your OpenRouter API Key", type="password")
st.sidebar.markdown("Get a key from [openrouter.ai](https://openrouter.ai)")

# Base URL override for OpenRouter
openai.api_key = api_key
openai.api_base = "https://openrouter.ai/api/v1"

# --- Choose Model ---
model = st.sidebar.selectbox("🤖 Choose AI Model", [
    "deepseek-ai/deepseek-coder",
    "mistralai/mixtral-8x7b-instruct",
    "meta-llama/llama-3-70b-instruct",
    "anthropic/claude-3-sonnet"
])

# --- File Upload ---
uploaded_file = st.file_uploader("Upload your file (PDF, DOCX, TXT, PPTX)", type=["pdf", "docx", "txt", "pptx"])

# --- File Handlers ---
def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_txt(file):
    return file.read().decode('utf-8')

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def download_as_pdf(text, filename="summary_output.pdf"):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for line in text.split("\n"):
        pdf.multi_cell(0, 10, line)
    output = io.BytesIO()
    pdf.output(output)
    return output

# --- Main App ---
if uploaded_file and api_key:
    file_ext = uploaded_file.name.split(".")[-1]

    with st.spinner("🔍 Extracting text..."):
        if file_ext == "pdf":
            text = extract_text_from_pdf(uploaded_file)
        elif file_ext == "docx":
            text = extract_text_from_docx(uploaded_file)
        elif file_ext == "txt":
            text = extract_text_from_txt(uploaded_file)
        elif file_ext == "pptx":
            text = extract_text_from_pptx(uploaded_file)
        else:
            st.error("Unsupported file type.")
            st.stop()

    # Prompt to summarize and visualize
    prompt = f"""
    Please summarize the following document in 5–7 lines.
    Then give bullet-point notes by topic.
    Then create a process flowchart in Mermaid.js markdown.

    Content:
    {text}
    """

    with st.spinner("🤖 Thinking with your chosen model..."):
        try:
            response = openai.ChatCompletion.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.4
            )
            result = response.choices[0].message["content"]
            st.subheader("📝 Summary, Notes, and Flowchart")
            st.text_area("Result", result, height=400)

            pdf_data = download_as_pdf(result)
            st.download_button("📥 Download as PDF", data=pdf_data, file_name="summary_output.pdf", mime="application/pdf")

            st.success("✅ Done!")
        except Exception as e:
            st.error(f"🚫 Error: {str(e)}")
